from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Helper: add styled slide with title & content
def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(10, 10, 40)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 215, 0)

    body_shape = slide.placeholders[1]
    body_shape.text = content
    for paragraph in body_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(230, 230, 250)
        paragraph.alignment = PP_ALIGN.LEFT

# Slide 1 - Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(5, 5, 30)
txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Astral Draw"
p.font.size = Pt(68)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 215, 0)
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "The Cosmic NFT Lottery on Hedera"
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(200, 200, 255)
p.alignment = PP_ALIGN.CENTER

# Executive Summary
add_slide("Executive Summary",
"Astral Draw is a decentralized, transparent, and gamified NFT-powered lottery platform built on Hedera Hashgraph. "
"It merges blockchain fairness with cosmic storytelling, allowing players to forge Star Keys (NFTs) and participate in Astral Convergences (lottery events). "
"The mission: redefine luck through transparency, governance, and innovation.")

# Problem
add_slide("The Problem",
"• Traditional lotteries lack transparency and fairness.\n"
"• Centralized control reduces trust and engagement.\n"
"• Younger digital audiences are disengaged from legacy lottery systems.\n"
"• Web3 entertainment lacks story-driven and accessible lottery mechanics.")

# Opportunity
add_slide("The Opportunity",
"• $400B global lottery market ripe for decentralization.\n"
"• Web3 gaming and NFT ecosystems exceed $50B.\n"
"• Astral Draw bridges these worlds with gamified fairness and on-chain proof.\n"
"• Early adoption in Hedera ecosystem positions us ahead of competitors.")

# Solution
add_slide("The Solution",
"• Blockchain-powered cosmic lottery where every draw is verifiable.\n"
"• NFT Star Keys unlock participation and rewards.\n"
"• Hedera Consensus Service ensures real-time fairness.\n"
"• DAO governance allows players to co-own the platform.")

# Product Overview
add_slide("Product Overview",
"• NFT-Gated Lottery Participation (Star Keys)\n"
"• DAO Governance via Hedera Consensus Service\n"
"• Smart Contract-Based Prize Distribution\n"
"• NFT Marketplace & Player Dashboard\n"
"• Interactive Cosmic Visual Experience")

# Technology Stack
add_slide("Technology Stack",
"• Hedera Hashgraph (HCS, HTS, HFS)\n"
"• HIP-412 NFT minting standard\n"
"• IPFS decentralized metadata\n"
"• Django + React + Web3.js for front-end and backend\n"
"• Blade & HashPack Wallet integrations")

# Competitive Advantage
add_slide("Competitive Advantage",
"• Fully transparent draws (Hedera)\n"
"• NFT-based entry ensures real ownership\n"
"• Lower transaction fees than Ethereum-based systems\n"
"• Cosmic branding & immersive storytelling\n"
"• Real utility beyond speculation")

# Tokenomics
add_slide("Tokenomics",
"• Token Name: ASTRAL\n"
"• Total Supply: 100M\n"
"• Allocation: 30% community, 20% team, 20% rewards, 15% liquidity, 15% investors\n"
"• Utility: Governance, staking, NFT minting, prize entry\n"
"• Deflationary mechanics via burn events")

# Business Model
add_slide("Business Model",
"• NFT Star Key sales\n"
"• 3% commission from prize pools\n"
"• NFT trading marketplace fees\n"
"• Sponsored Astral Events\n"
"• Premium membership tiers for players")

# Go-To-Market
add_slide("Go-To-Market Strategy",
"• Phase 1: Web3 community partnerships & testnet launch\n"
"• Phase 2: Influencer campaigns, NFT drops, and social tournaments\n"
"• Phase 3: Global marketing & DAO expansion\n"
"• Long-term: Integration with entertainment & esports sectors")

# Traction
add_slide("Traction & Milestones",
"• Prototype operational on Hedera Testnet\n"
"• Smart contract lottery logic implemented\n"
"• Early 100+ community members testing beta\n"
"• Demo video published and reviewed by experts")

# Governance
add_slide("Governance & Community",
"• DAO governance through HCS proposals\n"
"• NFT holders vote on event themes and prize pools\n"
"• Transparent voting and real-time results\n"
"• Incentives for participation and community growth")

# Financials
add_slide("Financial Projections (3 Years)",
"• Year 1: $50K NFT and event revenue\n"
"• Year 2: $250K from partnerships & sponsorships\n"
"• Year 3: $1M+ through global scaling\n"
"• Long-term: Expansion into Web3 entertainment & cross-chain NFTs")

# The Ask
add_slide("The Ask",
"Seeking $100,000 seed investment to:\n"
"• Scale infrastructure and marketing\n"
"• Expand smart contract ecosystem\n"
"• Onboard partners and artists\n"
"• Launch mainnet version with DAO activation")

# Team
add_slide("The Team",
"• Titos Kipkoech – Founder & Lead Developer\n"
"• Blockchain Researchers – Smart Contract Development\n"
"• UI/UX Designers – Cosmic Game Experience\n"
"• Advisors – Web3 Strategy and Compliance")

# Impact
add_slide("Impact",
"• Promotes fairness and transparency in gaming.\n"
"• Fosters digital inclusion and ownership.\n"
"• Supports creative economies through NFT rewards.\n"
"• Aligns with SDGs for innovation and equality.")

# Vision
add_slide("Vision & Closing",
"\"Astral Draw: Where Destiny Meets Technology.\"\n\n"
"Join us in reimagining global lotteries — transparent, fair, and cosmic.\n"
"Powered by Hedera. Guided by the stars.")

# Save final branded deck
path = "Astral_Draw_new_Pitch_Deck.pptx"
prs.save(path)

path
